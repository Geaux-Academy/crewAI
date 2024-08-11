from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
presentation = Presentation()

# Slide 1 - Title Slide
slide_1 = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]

title.text = "Proactive Steps CrewAI Assistant GPT Should Have Taken"
subtitle.text = "A 30-Minute Lecture on Best Practices"

# Slide 2 - Introduction
slide_2 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_2.shapes.title
content = slide_2.shapes.placeholders[1]

title.text = "Introduction"
content.text = (
    "In this lecture, we will discuss how CrewAI Assistant GPT should have "
    "proactively handled the project setup to avoid issues and ensure success. "
    "We'll cover key areas including understanding project requirements, environment setup, "
    "code structure, and error handling."
)

# Slide 3 - Understanding the Project Requirements
slide_3 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_3.shapes.title
content = slide_3.shapes.placeholders[1]

title.text = "Understanding the Project Requirements"
content.text = (
    "1. Clarify Objectives: Ask probing questions to fully understand the project's goals.\n"
    "2. Outline the Project Plan: Define agents, tasks, tools, and processes before implementation.\n"
    "3. Confirm Plan with User: Ensure the user agrees with the outlined plan before proceeding."
)

# Slide 4 - Proactive Environment Setup
slide_4 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_4.shapes.title
content = slide_4.shapes.placeholders[1]

title.text = "Proactive Environment Setup"
content.text = (
    "1. Create a Virtual Environment: Isolate dependencies to avoid conflicts.\n"
    "2. Install Required Packages: Ensure all necessary packages are installed.\n"
    "3. Handle Environment Variables: Use a .env file and python-dotenv to manage API keys.\n"
    "4. Configure PYTHONPATH: Set PYTHONPATH to include the src directory."
)

# Slide 5 - Code Structure and Import Management
slide_5 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_5.shapes.title
content = slide_5.shapes.placeholders[1]

title.text = "Code Structure and Import Management"
content.text = (
    "1. Use Absolute Imports: Avoid relative imports for better reliability.\n"
    "2. Verify Directory Structure: Ensure directories have __init__.py files.\n"
    "3. Simplify Imports: Keep imports straightforward to avoid ModuleNotFoundError."
)

# Slide 6 - Anticipating and Handling Errors
slide_6 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_6.shapes.title
content = slide_6.shapes.placeholders[1]

title.text = "Anticipating and Handling Errors"
content.text = (
    "1. Preemptive Debugging: Add logging to verify configurations.\n"
    "2. Error Handling Strategies: Discuss common errors and how to handle them.\n"
    "3. Incremental Testing: Test each component step-by-step."
)

# Slide 7 - Conclusion
slide_7 = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide_7.shapes.title
content = slide_7.shapes.placeholders[1]

title.text = "Conclusion"
content.text = (
    "By following these proactive steps, CrewAI Assistant GPT can avoid common pitfalls "
    "and ensure successful project setups. Clarify requirements, set up environments correctly, "
    "manage code structure, and anticipate errors to enhance efficiency and reliability."
)

# Save the presentation
presentation.save("CrewAI_Assistant_Proactive_Steps_Lecture.pptx")
